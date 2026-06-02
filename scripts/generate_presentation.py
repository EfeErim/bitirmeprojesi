#!/usr/bin/env python3
"""Convert SVGs in docs/presentation to slide-ready PNGs and build a PPTX.

Produces PNGs into `docs/presentation/pngs/` and `docs/presentation/presentation.pptx`.
Requires: pillow, python-pptx, and one of cairosvg, Inkscape, Edge, or Chrome.
"""
import shutil
import subprocess
from io import BytesIO
from pathlib import Path

from PIL import Image

try:
    import cairosvg
    HAS_CAIROSVG = True
except Exception:
    HAS_CAIROSVG = False
from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parents[1]
PRES_DIR = ROOT / "docs" / "presentation"
OUT_DIR = PRES_DIR / "pngs"
OUT_DIR.mkdir(parents=True, exist_ok=True)

def svg_to_png_centered(svg_path: Path, out_path: Path, canvas=(1920,1080), margin=80):
    # Render SVG to PNG bytes via cairosvg, Inkscape, or a Chromium browser.
    img = None
    if HAS_CAIROSVG:
        png_bytes = cairosvg.svg2png(url=str(svg_path))
        img = Image.open(BytesIO(png_bytes)).convert("RGBA")
    else:
        # try inkscape from PATH or common install location
        inkscape_cmd = shutil.which("inkscape")
        if not inkscape_cmd:
            # try program files path (common on Windows)
            possible = Path(r"C:/Program Files/Inkscape/bin/inkscape.com")
            if possible.exists():
                inkscape_cmd = str(possible)

        if inkscape_cmd:
            tmp_png = out_path.with_suffix('.tmp.png')
            cmd = [inkscape_cmd, str(svg_path), "--export-type=png", f"--export-filename={tmp_png}"]
            subprocess.run(cmd, check=True)
            img = Image.open(tmp_png).convert("RGBA")
            try:
                tmp_png.unlink()
            except Exception:
                pass
        else:
            browser_cmd = next(
                (
                    command
                    for command in (
                        shutil.which("msedge"),
                        shutil.which("chrome"),
                        r"C:/Program Files (x86)/Microsoft/Edge/Application/msedge.exe",
                        r"C:/Program Files/Microsoft/Edge/Application/msedge.exe",
                        r"C:/Program Files/Google/Chrome/Application/chrome.exe",
                    )
                    if command and Path(command).exists()
                ),
                None,
            )
            if not browser_cmd:
                raise RuntimeError("No renderer available: install cairosvg, Inkscape, Edge, or Chrome")
            tmp_png = out_path.with_suffix(".tmp.png")
            cmd = [
                str(browser_cmd),
                "--headless",
                "--disable-gpu",
                "--hide-scrollbars",
                f"--window-size={canvas[0]},{canvas[1]}",
                f"--screenshot={tmp_png}",
                svg_path.resolve().as_uri(),
            ]
            subprocess.run(cmd, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            img = Image.open(tmp_png).convert("RGBA")
            try:
                tmp_png.unlink()
            except Exception:
                pass

    canvas_w, canvas_h = canvas
    bg = Image.new("RGBA", canvas, (255,255,255,255))

    # compute target box inside margins
    max_w = canvas_w - margin*2
    max_h = canvas_h - margin*2
    w, h = img.size
    scale = min(max_w / w, max_h / h, 1.0)
    new_size = (int(w*scale), int(h*scale))
    img = img.resize(new_size, Image.LANCZOS)

    paste_x = (canvas_w - img.width)//2
    paste_y = (canvas_h - img.height)//2
    bg.paste(img, (paste_x, paste_y), img)
    bg.convert("RGB").save(out_path, "PNG", quality=95)

def build_pptx(image_paths, out_pptx: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 at 96 DPI
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Model Architecture & Pipeline"
    subtitle.text = "Diagrams exported from repository"

    for img_path in image_paths:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = img_path.stem.replace('_',' ').title()

        # add picture centered with max width/height
        pic_w = prs.slide_width - Inches(1)
        pic_h = prs.slide_height - Inches(1.5)
        left = (prs.slide_width - pic_w) / 2
        top = (prs.slide_height - pic_h) / 2
        slide.shapes.add_picture(str(img_path), left, top, width=pic_w)

    prs.save(out_pptx)

def main():
    svgs = list(PRES_DIR.glob("*.svg"))
    if not svgs:
        print("No SVGs found in", PRES_DIR)
        return

    pngs = []
    for svg in svgs:
        out_png = OUT_DIR / (svg.stem + ".png")
        print("Rendering", svg, "->", out_png)
        try:
            svg_to_png_centered(svg, out_png)
            pngs.append(out_png)
        except Exception as e:
            print("Failed to render", svg, e)

    out_pptx = PRES_DIR / "presentation.pptx"
    print("Building PPTX ->", out_pptx)
    build_pptx(pngs, out_pptx)
    print("Done. PNGs in:", OUT_DIR)

if __name__ == '__main__':
    main()
